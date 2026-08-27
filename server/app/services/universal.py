from __future__ import annotations

import csv
import html
import io
import shutil
import subprocess
import tempfile
from pathlib import Path

from PIL import Image

from app.core.config import settings
from app.core.files import command_path, run_command
from app.core.formats import (
    AUDIO_EXTENSIONS,
    IMAGE_EXTENSIONS,
    MEDIA_IMAGE_TARGETS,
    RENDER_IMAGE_TARGETS,
    PANDOC_INPUTS,
    PANDOC_OUTPUTS,
    PRESENTATIONS,
    SPREADSHEETS,
    TEXT_DOCUMENTS,
    VIDEO_EXTENSIONS,
    WORD_DOCUMENTS,
    extension_of,
    tool_state,
)
from app.services import archives, images, pdf_tools


def _soffice() -> str:
    return command_path(
        settings.soffice_bin,
        [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ],
    ) or ""


def _pandoc() -> str:
    return command_path(
        "pandoc",
        [r"C:\Program Files\Pandoc\pandoc.exe"],
    ) or ""


def _ffmpeg() -> str:
    return command_path(
        "ffmpeg",
        [
            r"C:\Program Files\ffmpeg\bin\ffmpeg.exe",
            r"C:\ffmpeg\bin\ffmpeg.exe",
        ],
    ) or ""


def _text_from_pdf(input_path: Path) -> str:
    import fitz

    doc = fitz.open(input_path)
    try:
        pages: list[str] = []
        for index, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            pages.append(f"## Page {index}\n\n{text}" if text else f"## Page {index}\n")
        return "\n\n".join(pages)
    finally:
        doc.close()


def _write_text_output(text: str, target_ext: str, output_path: Path) -> Path:
    target_ext = target_ext.lower()
    if target_ext in {"txt", "md", "rst", "org", "tex"}:
        output_path.write_text(text, encoding="utf-8")
        return output_path
    if target_ext in {"html", "htm"}:
        body = "<br>\n".join(html.escape(line) for line in text.splitlines())
        output_path.write_text(
            "<!doctype html><html><head><meta charset='utf-8'></head><body><p>"
            + body
            + "</p></body></html>",
            encoding="utf-8",
        )
        return output_path
    raise ValueError(f"No direct text writer for .{target_ext}")


def _pandoc_convert(input_path: Path, target_ext: str, out_dir: Path, source_ext: str | None = None) -> Path:
    pandoc = _pandoc()
    if not pandoc:
        raise RuntimeError("Pandoc is required for this conversion. Install Pandoc and restart Kiwi.")

    source_ext = source_ext or extension_of(input_path.name)
    from_fmt = PANDOC_INPUTS.get(source_ext)
    to_fmt = PANDOC_OUTPUTS.get(target_ext)
    if not from_fmt or not to_fmt:
        raise ValueError(f"Pandoc does not provide a reader/writer for .{source_ext} -> .{target_ext}.")

    output_path = out_dir / f"{input_path.stem}.{target_ext}"
    args = [
        pandoc,
        "--sandbox",
        "--from", from_fmt,
        "--to", to_fmt,
        "--standalone",
        "--output", str(output_path),
        str(input_path),
    ]
    run_command(args, timeout=240)
    if not output_path.exists():
        raise RuntimeError("Pandoc did not produce the requested output file.")
    return output_path


def _office_convert(input_path: Path, target_ext: str, out_dir: Path) -> Path:
    soffice = _soffice()
    if not soffice:
        raise RuntimeError("LibreOffice is required for this conversion. Install LibreOffice and restart Kiwi.")

    output_path = out_dir / f"{input_path.stem}.{target_ext}"
    filter_name = {
        "doc": "doc:MS Word 97",
        "ppt": "ppt:MS PowerPoint 97",
        "xls": "xls:MS Excel 97",
        "csv": "csv",
        "html": "html",
        "htm": "html",
    }.get(target_ext, target_ext)
    args = [
        soffice,
        "--headless",
        "--norestore",
        "--nolockcheck",
        "--nodefault",
        "--convert-to", filter_name,
        "--outdir", str(out_dir),
        str(input_path),
    ]
    run_command(args, timeout=240)
    if output_path.exists():
        return output_path

    candidates = [p for p in out_dir.glob(f"{input_path.stem}.*") if p.is_file()]
    if not candidates:
        raise RuntimeError("LibreOffice completed without producing an output file.")
    return candidates[0]


def _office_to_images(input_path: Path, target_ext: str, out_dir: Path) -> list[Path]:
    with tempfile.TemporaryDirectory(prefix="kiwi-office-") as tmp:
        tmp_dir = Path(tmp)
        pdf = _office_convert(input_path, "pdf", tmp_dir)
        rendered = pdf_tools.pdf_to_images(pdf, out_dir, target_ext, dpi=150)
        return rendered


def _pdf_tables(input_path: Path, target_ext: str, out_dir: Path) -> Path:
    try:
        import pdfplumber
    except ImportError as exc:
        raise RuntimeError("PDF table extraction requires pdfplumber. Reinstall Kiwi dependencies.") from exc

    rows: list[list[str]] = []
    with pdfplumber.open(input_path) as pdf:
        for page in pdf.pages:
            for table in page.extract_tables() or []:
                for row in table:
                    cleaned = [str(cell or "").strip() for cell in row]
                    if any(cleaned):
                        rows.append(cleaned)

    if not rows:
        raise RuntimeError(
            "No tables could be extracted from this PDF. A general PDF does not necessarily contain spreadsheet-shaped data."
        )

    if target_ext == "csv":
        output = out_dir / f"{input_path.stem}.csv"
        with output.open("w", newline="", encoding="utf-8-sig") as fp:
            writer = csv.writer(fp)
            writer.writerows(rows)
        return output

    if target_ext in {"xlsx", "xls", "ods"}:
        if target_ext == "xlsx":
            from openpyxl import Workbook

            output = out_dir / f"{input_path.stem}.xlsx"
            wb = Workbook()
            ws = wb.active
            for row in rows:
                ws.append(row)
            wb.save(output)
            return output

        csv_path = out_dir / f"{input_path.stem}.csv"
        with csv_path.open("w", newline="", encoding="utf-8-sig") as fp:
            csv.writer(fp).writerows(rows)
        return _office_convert(csv_path, target_ext, out_dir)

    raise ValueError(f"Unsupported table output .{target_ext}")


def _image_to_document(input_path: Path, target_ext: str, out_dir: Path) -> Path:
    image = Image.open(input_path).convert("RGB")
    try:
        if target_ext == "pdf":
            output = out_dir / f"{input_path.stem}.pdf"
            image.save(output, "PDF", resolution=150.0)
            return output

        if target_ext in {"docx", "doc", "odt", "rtf"}:
            from docx import Document
            from docx.shared import Inches

            temp_docx = out_dir / f"{input_path.stem}.docx"
            doc = Document()
            section = doc.sections[0]
            section.top_margin = Inches(0.5)
            section.bottom_margin = Inches(0.5)
            section.left_margin = Inches(0.5)
            section.right_margin = Inches(0.5)
            paragraph = doc.add_paragraph()
            paragraph.alignment = 1
            image.save(out_dir / f"{input_path.stem}_page.png", "PNG")
            paragraph.add_run().add_picture(str(out_dir / f"{input_path.stem}_page.png"), width=Inches(7.0))
            doc.save(temp_docx)
            (out_dir / f"{input_path.stem}_page.png").unlink(missing_ok=True)
            if target_ext == "docx":
                return temp_docx
            return _office_convert(temp_docx, target_ext, out_dir)

        if target_ext in {"pptx", "ppt", "odp"}:
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation()
            while len(prs.slides):
                r_id = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(r_id)
                del prs.slides._sldIdLst[0]
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            tmp = out_dir / f"{input_path.stem}_slide.png"
            image.save(tmp, "PNG")
            slide.shapes.add_picture(str(tmp), 0, 0, width=prs.slide_width, height=prs.slide_height)
            pptx_output = out_dir / f"{input_path.stem}.pptx"
            prs.save(pptx_output)
            tmp.unlink(missing_ok=True)
            if target_ext == "pptx":
                return pptx_output
            return _office_convert(pptx_output, target_ext, out_dir)

        if target_ext in {"txt", "md"}:
            import pytesseract
            tesseract = command_path(settings.tesseract_bin, [r"C:\Program Files\Tesseract-OCR\tesseract.exe"])
            if not tesseract:
                raise RuntimeError("OCR to text requires Tesseract. Install Tesseract and restart Kiwi.")
            pytesseract.pytesseract.tesseract_cmd = tesseract
            text = pytesseract.image_to_string(image)
            output = out_dir / f"{input_path.stem}.{target_ext}"
            output.write_text(text, encoding="utf-8")
            return output

        if target_ext in {"html", "htm"}:
            import base64
            buf = io.BytesIO()
            image.save(buf, "PNG")
            encoded = base64.b64encode(buf.getvalue()).decode("ascii")
            output = out_dir / f"{input_path.stem}.{target_ext}"
            output.write_text(
                f"<!doctype html><html><body><img src='data:image/png;base64,{encoded}' alt='{html.escape(input_path.stem)}'></body></html>",
                encoding="utf-8",
            )
            return output
    finally:
        image.close()

    raise ValueError(f"Cannot convert image to .{target_ext}")


def _text_to_image(input_path: Path, target_ext: str, out_dir: Path) -> Path:
    from PIL import ImageDraw, ImageFont

    text = input_path.read_text(encoding="utf-8", errors="replace")
    font_candidates = [
        r"C:\Windows\Fonts\arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    font = None
    for candidate in font_candidates:
        if Path(candidate).exists():
            try:
                font = ImageFont.truetype(candidate, 28)
                break
            except OSError:
                pass
    if font is None:
        font = ImageFont.load_default()

    lines = text.splitlines() or [""]
    line_h = 38
    width = 1400
    height = max(200, min(12000, 60 + len(lines) * line_h))
    image = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(image)

    y = 30
    for line in lines:
        draw.text((36, y), line[:140], fill="black", font=font)
        y += line_h
        if y >= height - line_h:
            break

    output = out_dir / f"{input_path.stem}.{target_ext}"
    fmt = "JPEG" if target_ext in {"jpg", "jpeg"} else target_ext.upper()
    image.save(output, format=fmt, quality=92 if fmt == "JPEG" else 95)
    image.close()
    return output


def _text_to_spreadsheet(input_path: Path, target_ext: str, out_dir: Path) -> Path:
    import csv as csv_module

    lines = input_path.read_text(encoding="utf-8", errors="replace").splitlines()
    rows: list[list[str]] = []

    # Markdown tables get parsed structurally. Plain text becomes one column.
    table_lines = [line.strip() for line in lines if "|" in line and line.strip().startswith("|")]
    if len(table_lines) >= 2:
        for line in table_lines:
            cells = [cell.strip() for cell in line.strip().strip("|").split("|")]
            if cells and all(set(cell) <= {"-", ":", " "} for cell in cells):
                continue
            rows.append(cells)
    else:
        rows = [[line] for line in lines if line.strip()] or [[""]]

    if target_ext == "csv":
        output = out_dir / f"{input_path.stem}.csv"
        with output.open("w", newline="", encoding="utf-8-sig") as fp:
            writer = csv_module.writer(fp)
            writer.writerows(rows)
        return output

    if target_ext == "xlsx":
        from openpyxl import Workbook
        output = out_dir / f"{input_path.stem}.xlsx"
        wb = Workbook()
        ws = wb.active
        for row in rows:
            ws.append(row)
        wb.save(output)
        return output

    if target_ext in {"xls", "ods"}:
        csv_path = _text_to_spreadsheet(input_path, "csv", out_dir)
        return _office_convert(csv_path, target_ext, out_dir)

    raise ValueError(f"Unsupported spreadsheet output .{target_ext}")


def _media_to_image(input_path: Path, target_ext: str, out_dir: Path) -> Path:
    ffmpeg = _ffmpeg()
    if not ffmpeg:
        raise RuntimeError("Media image rendering requires FFmpeg. Install FFmpeg and restart Kiwi.")
    output = out_dir / f"{input_path.stem}.{target_ext}"
    source_ext = extension_of(input_path.name)
    if source_ext in AUDIO_EXTENSIONS:
        filter_name = "showwavespic=s=1600x900:colors=white"
        args = [ffmpeg, "-y", "-i", str(input_path), "-filter_complex", filter_name, "-frames:v", "1", str(output)]
    else:
        args = [ffmpeg, "-y", "-i", str(input_path), "-frames:v", "1", str(output)]
    run_command(args, timeout=300)
    if not output.exists():
        raise RuntimeError("FFmpeg did not produce the requested image.")
    if target_ext == "pdf":
        pdf = _image_to_document(output, "pdf", out_dir)
        output.unlink(missing_ok=True)
        return pdf
    return output


def _archive_convert(input_path: Path, target_ext: str, out_dir: Path) -> Path:
    extract_root = out_dir / "_repack"
    extracted = archives.extract_archive(input_path, extract_root)
    try:
        # extract_archive returns files plus no package marker now. Repack the files with safe relative names.
        return archives.create_archive(extracted, out_dir, target_ext, archives.archive_stem(input_path))
    finally:
        shutil.rmtree(extract_root, ignore_errors=True)


def _media_convert(input_path: Path, target_ext: str, out_dir: Path) -> Path:
    ffmpeg = _ffmpeg()
    if not ffmpeg:
        raise RuntimeError("Media conversion requires FFmpeg. Install FFmpeg and restart Kiwi.")

    output = out_dir / f"{input_path.stem}.{target_ext}"
    source_ext = extension_of(input_path.name)

    if source_ext in AUDIO_EXTENSIONS and target_ext in VIDEO_EXTENSIONS:
        args = [
            ffmpeg, "-y",
            "-f", "lavfi", "-i", "color=c=black:s=1280x720:r=30",
            "-i", str(input_path),
            "-shortest",
            "-map", "0:v:0", "-map", "1:a:0",
            "-c:v", "libx264", "-pix_fmt", "yuv420p",
            "-c:a", "aac",
            "-map_metadata", "0",
            str(output),
        ]
    elif source_ext in VIDEO_EXTENSIONS and target_ext in AUDIO_EXTENSIONS:
        args = [
            ffmpeg, "-y",
            "-i", str(input_path),
            "-vn",
            "-map", "0:a:0?",
            "-map_metadata", "0",
            str(output),
        ]
    else:
        args = [
            ffmpeg, "-y",
            "-i", str(input_path),
            "-map_metadata", "0",
            str(output),
        ]

    run_command(args, timeout=600)
    if not output.exists():
        raise RuntimeError("FFmpeg did not produce the requested output file.")
    return output


def convert_any(input_path: Path, target_ext: str, out_dir: Path) -> list[Path]:
    source_ext = extension_of(input_path.name)
    target_ext = target_ext.lower().lstrip(".")
    out_dir.mkdir(parents=True, exist_ok=True)

    if source_ext == target_ext:
        raise ValueError("Source and target formats are identical.")

    # PDF paths.
    if source_ext == "pdf":
        if target_ext == "txt":
            output = out_dir / f"{input_path.stem}.txt"
            return [_write_text_output(_text_from_pdf(input_path), "txt", output)]
        if target_ext == "html":
            output = out_dir / f"{input_path.stem}.html"
            return [_write_text_output(_text_from_pdf(input_path), "html", output)]
        if target_ext == "docx":
            return [documents_pdf_to_docx(input_path, out_dir)]
        if target_ext == "pptx":
            return [documents_pdf_to_pptx(input_path, out_dir)]
        if target_ext in RENDER_IMAGE_TARGETS:
            fmt = "jpg" if target_ext == "jpeg" else target_ext
            return [p for p in pdf_tools.pdf_to_images(input_path, out_dir, fmt, dpi=150)]
        if target_ext in {"csv", "xlsx", "xls", "ods"}:
            return [_pdf_tables(input_path, target_ext, out_dir)]
        if target_ext in {"doc", "odt", "rtf"}:
            intermediate = documents_pdf_to_docx(input_path, out_dir)
            if target_ext == "docx":
                return [intermediate]
            return [_office_convert(intermediate, target_ext, out_dir)]
        if target_ext in {"ppt", "odp"}:
            intermediate = documents_pdf_to_pptx(input_path, out_dir)
            return [_office_convert(intermediate, target_ext, out_dir)]
        if target_ext in {"md", "epub", "tex", "rst", "org"}:
            text_path = out_dir / f"{input_path.stem}.md"
            text_path.write_text(_text_from_pdf(input_path), encoding="utf-8")
            result = _pandoc_convert(text_path, target_ext, out_dir, "md")
            text_path.unlink(missing_ok=True)
            return [result]

    # Images.
    if source_ext in IMAGE_EXTENSIONS:
        if target_ext in IMAGE_EXTENSIONS:
            return [images.convert_image(input_path, target_ext, out_dir)]
        return [_image_to_document(input_path, target_ext, out_dir)]

    # Archive repackaging is useful when the container itself is the target.
    if source_ext in {"zip", "7z", "tar", "tar.gz", "tgz", "tar.bz2", "tbz2", "tar.xz", "txz"} and target_ext in {"zip", "7z", "tar"}:
        return [_archive_convert(input_path, target_ext, out_dir)]

    # Audio/video.
    if source_ext in AUDIO_EXTENSIONS or source_ext in VIDEO_EXTENSIONS:
        if target_ext in MEDIA_IMAGE_TARGETS:
            fmt = "jpg" if target_ext == "jpeg" else target_ext
            return [_media_to_image(input_path, fmt, out_dir)]
        if target_ext == "pdf":
            jpg = _media_to_image(input_path, "jpg", out_dir)
            try:
                return [_image_to_document(jpg, "pdf", out_dir)]
            finally:
                jpg.unlink(missing_ok=True)
        return [_media_convert(input_path, target_ext, out_dir)]

    # Text to images is a real rendered representation, not a fake extension rename.
    if source_ext in TEXT_DOCUMENTS and target_ext in RENDER_IMAGE_TARGETS:
        return [_text_to_image(input_path, target_ext, out_dir)]

    if source_ext in TEXT_DOCUMENTS and target_ext in {"csv", "xlsx", "xls", "ods"}:
        return [_text_to_spreadsheet(input_path, target_ext, out_dir)]

    # Broad text/code bridge. Unsupported text extensions are safely parsed as plain text.
    if source_ext in TEXT_DOCUMENTS:
        if target_ext in PANDOC_OUTPUTS:
            return [_pandoc_convert(input_path, target_ext, out_dir, PANDOC_INPUTS.get(source_ext, "plain"))]
        if target_ext == "pdf" and _soffice():
            docx = _pandoc_convert(input_path, "docx", out_dir, PANDOC_INPUTS.get(source_ext, "plain"))
            return [_office_convert(docx, "pdf", out_dir)]

    # Native Pandoc document/markup bridge.
    if source_ext in PANDOC_INPUTS and target_ext in PANDOC_OUTPUTS:
        return [_pandoc_convert(input_path, target_ext, out_dir, source_ext)]

    # LibreOffice handles many legacy/Office conversions.
    if target_ext in {"pdf", "doc", "docx", "odt", "rtf", "txt", "html", "htm", "xlsx", "xls", "ods", "csv", "ppt", "pptx", "odp"}:
        return [_office_convert(input_path, target_ext, out_dir)]

    if target_ext in RENDER_IMAGE_TARGETS:
        return [p for p in _office_to_images(input_path, target_ext, out_dir)]

    raise RuntimeError(f"Kiwi has no conversion engine for .{source_ext} -> .{target_ext}.")


def documents_pdf_to_docx(input_path: Path, out_dir: Path) -> Path:
    # Keep the existing high-quality PDF block extractor.
    from app.services.documents import _pdf_to_docx
    return _pdf_to_docx(input_path, out_dir / f"{input_path.stem}.docx")


def documents_pdf_to_pptx(input_path: Path, out_dir: Path) -> Path:
    from app.services.documents import _pdf_to_pptx
    return _pdf_to_pptx(input_path, out_dir / f"{input_path.stem}.pptx")
