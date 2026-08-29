"""Family-aware LibreOffice conversion engine.

Fixes the root cause behind the PPTX -> DOCX failure: the old code treated
every "convert to .X" request as "pass X (or a hand-picked legacy filter
string) to `soffice --convert-to`", which only works when source and target
belong to the same LibreOffice document family (Writer stays Writer,
Impress stays Impress, Calc stays Calc). `--convert-to` has no cross-family
transform -- there is no "Impress document -> save as Word" filter, no
matter which filter name is used, because Impress simply cannot produce a
Writer document. Verified directly against LibreOffice 24.2:

    soffice --convert-to "docx:MS Word 2007 XML" fees_26_27_first.pptx
    -> SfxBaseModel::impl_store ... failed (0xc10)   # every docx filter, same result

    soffice --convert-to pptx test.docx
    -> Error: no export filter for test.pptx found, aborting.

So this module distinguishes:

  * OfficeFilterRegistry  -- the *correct* modern filter name for every
    extension Kiwi claims to support (never a guessed "doc:MS Word 97" for
    an OOXML target), keyed by the document family it belongs to.
  * same-family conversion -- a direct, single `--convert-to` call
    (Writer->Writer, Impress->Impress, Calc->Calc, or any family -> PDF,
    which every family can export).
  * cross-family conversion -- no direct LibreOffice filter exists. Kiwi
    performs a semantic bridge instead (e.g. PPTX -> DOCX pulls slide text
    into a flowing Word document; DOCX -> PPTX puts paragraphs onto slides)
    rather than pretending `--convert-to` can do it.

Every produced file is validated (exists, correct extension, non-empty,
openable by an appropriate parser) before being reported as a success --
see `validate_office_output`.
"""
from __future__ import annotations

import shutil
import tempfile
import uuid
import zipfile
from dataclasses import dataclass
from pathlib import Path

from app.core.config import settings
from app.core.errors import ConversionDiagnostics, KiwiConversionError, OutputValidationError
from app.core.files import command_path

import subprocess


# ---------------------------------------------------------------------------
# Filter registry
# ---------------------------------------------------------------------------

WRITER = "writer"       # word-processing documents
IMPRESS = "impress"     # presentations
CALC = "calc"           # spreadsheets
DRAW = "draw"           # vector drawings (not currently a Kiwi target, listed for completeness)


@dataclass(frozen=True)
class OfficeFilter:
    extension: str
    filter_name: str
    family: str
    expected_mime: str | None = None


# Verified against LibreOffice 24.2 filter names (see
# https://help.libreoffice.org/latest/ast/text/shared/guide/convertfilters.html
# and `soffice --convert-to <ext>:<help>` output). Modern OOXML formats use
# their real 2007+ filter names, never the legacy "97" binary filters --
# using a legacy filter for a modern extension is exactly the bug this
# registry exists to prevent.
OFFICE_FILTERS: dict[str, OfficeFilter] = {
    # Writer family
    "doc": OfficeFilter("doc", "MS Word 97", WRITER, "application/msword"),
    "docx": OfficeFilter("docx", "MS Word 2007 XML", WRITER,
                          "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
    "docm": OfficeFilter("docm", "MS Word 2007 XML VBA", WRITER,
                          "application/vnd.ms-word.document.macroEnabled.12"),
    "dot": OfficeFilter("dot", "MS Word 97 Vorlage", WRITER, "application/msword"),
    "dotx": OfficeFilter("dotx", "MS Word 2007 XML Template", WRITER,
                          "application/vnd.openxmlformats-officedocument.wordprocessingml.template"),
    "odt": OfficeFilter("odt", "writer8", WRITER, "application/vnd.oasis.opendocument.text"),
    "rtf": OfficeFilter("rtf", "Rich Text Format", WRITER, "application/rtf"),

    # Impress family
    "ppt": OfficeFilter("ppt", "MS PowerPoint 97", IMPRESS, "application/vnd.ms-powerpoint"),
    "pptx": OfficeFilter("pptx", "Impress MS PowerPoint 2007 XML", IMPRESS,
                          "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    "pptm": OfficeFilter("pptm", "Impress MS PowerPoint 2007 XML VBA", IMPRESS,
                          "application/vnd.ms-powerpoint.presentation.macroEnabled.12"),
    "pot": OfficeFilter("pot", "MS PowerPoint 97 Vorlage", IMPRESS, "application/vnd.ms-powerpoint"),
    "potx": OfficeFilter("potx", "Impress MS PowerPoint 2007 XML Template", IMPRESS,
                          "application/vnd.openxmlformats-officedocument.presentationml.template"),
    "odp": OfficeFilter("odp", "impress8", IMPRESS, "application/vnd.oasis.opendocument.presentation"),

    # Calc family
    "xls": OfficeFilter("xls", "MS Excel 97", CALC, "application/vnd.ms-excel"),
    "xlsx": OfficeFilter("xlsx", "Calc MS Excel 2007 XML", CALC,
                          "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
    "xlsm": OfficeFilter("xlsm", "Calc MS Excel 2007 VBA XML", CALC,
                          "application/vnd.ms-excel.sheet.macroEnabled.12"),
    "ods": OfficeFilter("ods", "calc8", CALC, "application/vnd.oasis.opendocument.spreadsheet"),
    "csv": OfficeFilter("csv", "Text - txt - csv (StarCalc)", CALC, "text/csv"),

    # Cross-family / render targets every family can export to
    "pdf": OfficeFilter("pdf", "writer_pdf_Export", WRITER, "application/pdf"),  # family overridden per-source below
    "html": OfficeFilter("html", "HTML (StarWriter)", WRITER, "text/html"),
    "htm": OfficeFilter("htm", "HTML (StarWriter)", WRITER, "text/html"),
}

# PDF/HTML export filter names differ per source family even though the
# target extension is the same.
_PDF_EXPORT_FILTER = {WRITER: "writer_pdf_Export", IMPRESS: "impress_pdf_Export", CALC: "calc_pdf_Export"}
_HTML_EXPORT_FILTER = {WRITER: "HTML (StarWriter)", IMPRESS: "impress_html_Export", CALC: "HTML (StarCalc)"}

FAMILY_OF_EXTENSION: dict[str, str] = {ext: f.family for ext, f in OFFICE_FILTERS.items()}
FAMILY_OF_EXTENSION.update({"docm": WRITER, "dot": WRITER, "dotx": WRITER})


def family_of(ext: str) -> str | None:
    return FAMILY_OF_EXTENSION.get(ext.lower().lstrip("."))


def resolve_filter(source_ext: str, target_ext: str) -> str:
    """Return the correct `--convert-to` filter token ("ext:Filter Name")
    for a same-family (or any-family -> pdf/html) conversion. Raises if the
    pair is cross-family and has no direct LibreOffice filter -- callers
    must use the semantic bridge instead.
    """
    source_ext = source_ext.lower().lstrip(".")
    target_ext = target_ext.lower().lstrip(".")
    source_family = family_of(source_ext)

    if target_ext == "pdf" and source_family in _PDF_EXPORT_FILTER:
        return f"pdf:{_PDF_EXPORT_FILTER[source_family]}"
    if target_ext in {"html", "htm"} and source_family in _HTML_EXPORT_FILTER:
        return f"{target_ext}:{_HTML_EXPORT_FILTER[source_family]}"

    target = OFFICE_FILTERS.get(target_ext)
    if target is None:
        raise KiwiConversionError(f"No LibreOffice filter is registered for .{target_ext}.")

    if target_ext == "csv":
        return f"csv:{target.filter_name}"

    if source_family is not None and source_family != target.family:
        raise KiwiConversionError(
            f"LibreOffice cannot directly convert .{source_ext} to .{target_ext} "
            f"({source_family} document family to {target.family} document family)."
        )
    return f"{target_ext}:{target.filter_name}"


def can_convert_directly(source_ext: str, target_ext: str) -> bool:
    """True if a direct same-family (or any-family -> pdf/html) LibreOffice
    filter exists for this pair."""
    source_ext = source_ext.lower().lstrip(".")
    target_ext = target_ext.lower().lstrip(".")
    if target_ext in {"pdf", "html", "htm"}:
        return family_of(source_ext) in (_PDF_EXPORT_FILTER if target_ext == "pdf" else _HTML_EXPORT_FILTER)
    target = OFFICE_FILTERS.get(target_ext)
    if target is None:
        return False
    source_family = family_of(source_ext)
    return source_family is not None and source_family == target.family


# ---------------------------------------------------------------------------
# Output validation (spec section 3 / 5): never trust "soffice exited 0" or
# "a file with a matching stem exists" as proof of success.
# ---------------------------------------------------------------------------

_ZIP_MAGIC = b"PK\x03\x04"
_OLE_MAGIC = b"\xd0\xcf\x11\xe0"  # legacy .doc/.ppt/.xls compound file


def _container_signature_ok(path: Path, ext: str) -> bool:
    ext = ext.lower()
    try:
        with path.open("rb") as fh:
            head = fh.read(8)
    except OSError:
        return False
    if ext in {"docx", "pptx", "xlsx", "docm", "pptm", "xlsm", "dotx", "potx", "odt", "odp", "ods"}:
        return head.startswith(_ZIP_MAGIC)
    if ext in {"doc", "ppt", "xls", "dot", "pot"}:
        return head.startswith(_OLE_MAGIC)
    if ext == "pdf":
        return head.startswith(b"%PDF")
    return True  # csv/html/rtf/txt: no reliable binary magic, size check below is enough


def _openable_by_parser(path: Path, ext: str) -> bool:
    ext = ext.lower()
    try:
        if ext == "docx":
            from docx import Document
            Document(str(path))
        elif ext == "pptx":
            from pptx import Presentation
            Presentation(str(path))
        elif ext == "xlsx":
            from openpyxl import load_workbook
            load_workbook(str(path), read_only=True).close()
        elif ext in {"docx", "pptx", "xlsx", "odt", "odp", "ods"}:
            with zipfile.ZipFile(path) as zf:
                if zf.testzip() is not None:
                    return False
        return True
    except Exception:
        return False


def validate_office_output(path: Path, expected_ext: str, *, engine: str = "libreoffice") -> Path:
    """Validate section 3's checklist. Raises OutputValidationError instead
    of letting a wrong-format or corrupt file silently pass as a success.
    """
    expected_ext = expected_ext.lower().lstrip(".")

    if not path.exists() or not path.is_file():
        raise OutputValidationError(
            f"The conversion tool did not produce the requested .{expected_ext} file.",
            ConversionDiagnostics(engine=engine, target=expected_ext, extra={"expected_path": str(path)}),
        )

    actual_ext = path.suffix.lower().lstrip(".")
    if actual_ext != expected_ext:
        raise OutputValidationError(
            f"The conversion tool produced a .{actual_ext} file instead of the requested .{expected_ext}.",
            ConversionDiagnostics(engine=engine, target=expected_ext, extra={"actual_path": str(path)}),
        )

    if path.stat().st_size == 0:
        raise OutputValidationError(
            f"The conversion produced an empty .{expected_ext} file.",
            ConversionDiagnostics(engine=engine, target=expected_ext),
        )

    if not _container_signature_ok(path, expected_ext):
        raise OutputValidationError(
            f"The .{expected_ext} file that was produced does not have a valid {expected_ext.upper()} container.",
            ConversionDiagnostics(engine=engine, target=expected_ext),
        )

    if not _openable_by_parser(path, expected_ext):
        raise OutputValidationError(
            f"The .{expected_ext} file that was produced could not be opened/parsed after conversion.",
            ConversionDiagnostics(engine=engine, target=expected_ext),
        )

    return path


# ---------------------------------------------------------------------------
# LibreOffice invocation, with isolated per-job user profiles (section 41).
# ---------------------------------------------------------------------------

def _soffice_bin() -> str:
    found = command_path(
        settings.soffice_bin if hasattr(settings, "soffice_bin") else "soffice",
        [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ],
    )
    if not found:
        raise KiwiConversionError("LibreOffice is required for this conversion. Install LibreOffice and restart Kiwi.")
    return found


def run_soffice_convert(input_path: Path, filter_token: str, out_dir: Path, *, timeout: int = 240) -> subprocess.CompletedProcess:
    """Run `soffice --convert-to <filter_token>` inside an isolated
    temporary user-profile directory, so concurrent conversions never share
    LibreOffice's lock/profile state (a shared profile is not
    concurrency-safe even with --headless -- simultaneous invocations can
    collide on the user installation lock).
    """
    binary = _soffice_bin()
    out_dir.mkdir(parents=True, exist_ok=True)
    profile_dir = Path(tempfile.gettempdir()) / f"kiwi-lo-profile-{uuid.uuid4().hex}"
    profile_dir.mkdir(parents=True, exist_ok=True)
    profile_uri = profile_dir.as_uri()

    args = [
        binary,
        "--headless",
        "--invisible",
        "--norestore",
        "--nolockcheck",
        "--nodefault",
        "--nofirststartwizard",
        f"-env:UserInstallation={profile_uri}",
        "--convert-to", filter_token,
        "--outdir", str(out_dir),
        str(input_path),
    ]
    try:
        return subprocess.run(args, capture_output=True, timeout=timeout, text=True)
    finally:
        shutil.rmtree(profile_dir, ignore_errors=True)


def office_convert_same_family(input_path: Path, target_ext: str, out_dir: Path, *, timeout: int = 240) -> Path:
    """Direct same-family (or any-family -> pdf/html) LibreOffice
    conversion, with a correctly-resolved filter and full output
    validation. Raises KiwiConversionError with safe messaging on failure
    -- never a bare CalledProcessError, never a silently-wrong file.
    """
    source_ext = input_path.suffix.lower().lstrip(".")
    target_ext = target_ext.lower().lstrip(".")
    filter_token = resolve_filter(source_ext, target_ext)

    result = run_soffice_convert(input_path, filter_token, out_dir, timeout=timeout)
    diagnostics = ConversionDiagnostics(
        engine="libreoffice",
        source=source_ext,
        target=target_ext,
        command=[Path(input_path).name, "--convert-to", filter_token],
        filter_name=filter_token,
        exit_code=result.returncode,
        stdout=result.stdout or "",
        stderr=result.stderr or "",
    )

    expected = out_dir / f"{input_path.stem}.{target_ext}"

    if result.returncode != 0:
        raise KiwiConversionError(
            f"LibreOffice could not convert .{source_ext} to .{target_ext}. The selected output filter failed.",
            diagnostics,
        )

    if not expected.exists():
        # Only look for files with the exact requested extension -- never
        # grab an arbitrary same-stem file (that was the old `candidates[0]`
        # bug: it could silently hand back a .odp when .docx was requested).
        matches = sorted(
            p for p in out_dir.glob(f"{input_path.stem}.*")
            if p.is_file() and p.suffix.lower().lstrip(".") == target_ext
        )
        if not matches:
            raise KiwiConversionError(
                f"LibreOffice completed without producing the requested .{target_ext} output.",
                diagnostics,
            )
        expected = matches[0]

    return validate_office_output(expected, target_ext, engine="libreoffice")


# ---------------------------------------------------------------------------
# Cross-family semantic bridges: PPTX <-> DOCX.
#
# There is no LibreOffice filter for these (verified above), because a
# presentation and a word-processing document are structurally different
# things. Instead of failing outright, Kiwi builds a lightweight
# DocumentIR-style intermediate (a flat list of blocks: heading/paragraph/
# table/image) and renders it into the target container. This is a
# "reconstruct", not a "render" (section 33): it preserves *content* and
# reading order, not the original visual layout, because a slide deck has
# no flowing-page layout to preserve in the first place.
# ---------------------------------------------------------------------------


@dataclass
class IRBlock:
    kind: str  # "heading" | "paragraph" | "table" | "image" | "page_break"
    text: str = ""
    level: int = 1
    rows: list[list[str]] | None = None
    image_path: str | None = None


def _pptx_to_ir(input_path: Path, tmp_dir: Path) -> list[IRBlock]:
    from pptx import Presentation

    prs = Presentation(str(input_path))
    blocks: list[IRBlock] = []
    img_counter = 0

    for slide_index, slide in enumerate(prs.slides):
        title_text = None
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            title_text = slide.shapes.title.text_frame.text.strip() or None
        blocks.append(IRBlock(kind="heading", text=title_text or f"Slide {slide_index + 1}", level=1))

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_text_frame and shape.text_frame.text.strip():
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs) or para.text
                    if text.strip():
                        blocks.append(IRBlock(kind="paragraph", text=text.strip()))
            elif shape.has_table:
                table = shape.table
                rows = [[cell.text for cell in row.cells] for row in table.rows]
                if rows:
                    blocks.append(IRBlock(kind="table", rows=rows))
            elif shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image = shape.image
                    img_counter += 1
                    img_path = tmp_dir / f"slide{slide_index}_img{img_counter}.{image.ext}"
                    img_path.write_bytes(image.blob)
                    blocks.append(IRBlock(kind="image", image_path=str(img_path)))
                except Exception:
                    continue

        # Speaker notes carry real authored content -- preserve them as a
        # labelled paragraph instead of silently dropping them.
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                blocks.append(IRBlock(kind="paragraph", text=f"[Speaker notes] {notes}"))

        blocks.append(IRBlock(kind="page_break"))

    return blocks


def _ir_to_docx(blocks: list[IRBlock], output_path: Path) -> Path:
    from docx import Document
    from docx.shared import Inches

    doc = Document()
    for block in blocks:
        if block.kind == "heading":
            doc.add_heading(block.text, level=min(max(block.level, 1), 4))
        elif block.kind == "paragraph":
            doc.add_paragraph(block.text)
        elif block.kind == "table" and block.rows:
            rows, cols = len(block.rows), max(len(r) for r in block.rows)
            table = doc.add_table(rows=rows, cols=cols)
            table.style = "Table Grid"
            for r, row in enumerate(block.rows):
                for c, cell_text in enumerate(row):
                    table.cell(r, c).text = cell_text
        elif block.kind == "image" and block.image_path and Path(block.image_path).exists():
            try:
                doc.add_picture(block.image_path, width=Inches(6.0))
            except Exception:
                pass
        elif block.kind == "page_break":
            doc.add_page_break()
    doc.save(output_path)
    return output_path


def _docx_to_ir(input_path: Path, tmp_dir: Path) -> list[IRBlock]:
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(str(input_path))
    blocks: list[IRBlock] = []

    def iter_block_items(document):
        # Walk body children in document order so tables interleave
        # correctly with paragraphs instead of all tables trailing at the
        # end (python-docx exposes .paragraphs and .tables as separate,
        # unordered lists).
        from docx.table import Table
        from docx.text.paragraph import Paragraph
        body = document.element.body
        for child in body.iterchildren():
            if child.tag == qn("w:p"):
                yield Paragraph(child, document)
            elif child.tag == qn("w:tbl"):
                yield Table(child, document)

    for item in iter_block_items(doc):
        cls_name = item.__class__.__name__
        if cls_name == "Paragraph":
            text = item.text.strip()
            if not text:
                continue
            style = (item.style.name or "").lower() if item.style else ""
            if "heading" in style or "title" in style:
                level = 1
                for ch in style:
                    if ch.isdigit():
                        level = int(ch)
                        break
                blocks.append(IRBlock(kind="heading", text=text, level=level))
            else:
                blocks.append(IRBlock(kind="paragraph", text=text))
        elif cls_name == "Table":
            rows = [[cell.text for cell in row.cells] for row in item.rows]
            if rows:
                blocks.append(IRBlock(kind="table", rows=rows))

    return blocks


def _ir_to_pptx(blocks: list[IRBlock], output_path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[1]  # Title and Content
    blank_layout = prs.slide_layouts[6]

    slide = None
    body_tf = None
    MAX_LINES_PER_SLIDE = 12
    lines_on_slide = 0

    def new_slide(title_text: str = ""):
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = title_text or " "
        tf = s.placeholders[1].text_frame
        tf.clear()
        return s, tf

    for block in blocks:
        if block.kind == "heading":
            slide, body_tf = new_slide(block.text)
            lines_on_slide = 0
        elif block.kind == "paragraph":
            if slide is None or lines_on_slide >= MAX_LINES_PER_SLIDE:
                slide, body_tf = new_slide("")
                lines_on_slide = 0
            p = body_tf.paragraphs[0] if lines_on_slide == 0 and not body_tf.paragraphs[0].runs else body_tf.add_paragraph()
            p.text = block.text
            p.font.size = Pt(18)
            lines_on_slide += 1
        elif block.kind == "table" and block.rows:
            rows, cols = len(block.rows), max(len(r) for r in block.rows)
            s = prs.slides.add_slide(blank_layout)
            table_shape = s.shapes.add_table(rows, cols, Inches(0.5), Inches(0.5), Inches(12), Inches(6))
            table = table_shape.table
            for r, row in enumerate(block.rows):
                for c in range(cols):
                    table.cell(r, c).text = row[c] if c < len(row) else ""
            slide, body_tf = None, None
        elif block.kind == "image" and block.image_path and Path(block.image_path).exists():
            s = prs.slides.add_slide(blank_layout)
            try:
                s.shapes.add_picture(block.image_path, Inches(1), Inches(1), width=Inches(10))
            except Exception:
                pass
            slide, body_tf = None, None
        # page_break: no-op for docx->pptx direction, slides are already discrete.

    if not prs.slides:
        new_slide("")

    prs.save(output_path)
    return output_path


def pptx_to_docx_semantic(input_path: Path, output_path: Path) -> Path:
    """PPTX -> DOCX via a slide-content extraction bridge (no LibreOffice
    filter exists for this cross-family pair -- see module docstring)."""
    with tempfile.TemporaryDirectory(prefix="kiwi-pptx2docx-") as tmp:
        blocks = _pptx_to_ir(input_path, Path(tmp))
        _ir_to_docx(blocks, output_path)
    return validate_office_output(output_path, "docx", engine="kiwi-semantic-bridge")


def docx_to_pptx_semantic(input_path: Path, output_path: Path) -> Path:
    """DOCX -> PPTX via a heading/paragraph-to-slide bridge (no LibreOffice
    filter exists for this cross-family pair -- see module docstring)."""
    with tempfile.TemporaryDirectory(prefix="kiwi-docx2pptx-") as tmp:
        blocks = _docx_to_ir(input_path, Path(tmp))
        _ir_to_pptx(blocks, output_path)
    return validate_office_output(output_path, "pptx", engine="kiwi-semantic-bridge")


# ---------------------------------------------------------------------------
# Public entry point used by universal.convert_any / documents.py.
# ---------------------------------------------------------------------------

# Cross-family pairs Kiwi has a real semantic bridge for. Anything else
# cross-family (e.g. xlsx -> pptx) has no bridge yet and should be reported
# as unsupported rather than silently attempted.
_SEMANTIC_BRIDGES = {
    ("pptx", "docx"): pptx_to_docx_semantic,
    ("docx", "pptx"): docx_to_pptx_semantic,
}


def office_convert(input_path: Path, target_ext: str, out_dir: Path, *, timeout: int = 240) -> Path:
    """Single entry point: same-family -> direct LibreOffice conversion;
    known cross-family pair -> semantic bridge; anything else -> a clear
    UnsupportedConversionError instead of a doomed `--convert-to` call.
    """
    from app.core.errors import UnsupportedConversionError

    source_ext = input_path.suffix.lower().lstrip(".")
    target_ext = target_ext.lower().lstrip(".")
    out_dir.mkdir(parents=True, exist_ok=True)
    output_path = out_dir / f"{input_path.stem}.{target_ext}"

    if can_convert_directly(source_ext, target_ext):
        return office_convert_same_family(input_path, target_ext, out_dir, timeout=timeout)

    bridge = _SEMANTIC_BRIDGES.get((source_ext, target_ext))
    if bridge is not None:
        return bridge(input_path, output_path)

    raise UnsupportedConversionError(
        f".{source_ext} to .{target_ext} has no direct LibreOffice filter and no semantic bridge in Kiwi yet."
    )
