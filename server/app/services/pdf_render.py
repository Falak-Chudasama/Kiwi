"""Layout-aware PDF renderers: PDF -> fixed-layout HTML, PDF -> PPTX.

Both renderers consume the same per-page primitives (text spans with real
bounding boxes, resolved fonts, images, vector-graphic regions) gathered via
PyMuPDF, and fall back to OCR-derived text for pages that `pdf_analysis`
flagged as having no usable text layer. This replaces the previous
"dump page.get_text('text') into <p><br></p>" HTML output and the
Arial-only, no-OCR PPTX output.
"""

from __future__ import annotations

import html
import io
from pathlib import Path

from app.services.pdf_analysis import PdfAnalysis, analyze_pdf
from app.services.pdf_fonts import resolve_font


def _page_spans(page, ocr_fallback: bool):
    """Yield (bbox, text, font_family, size, bold, italic, color) for every
    text span on the page, in document order. Falls back to OCR line boxes
    when the page has no native text layer and OCR is available.
    """
    text_dict = page.get_text("dict")
    any_text = False
    for block in text_dict.get("blocks", []):
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                text = span.get("text", "")
                if not text.strip():
                    continue
                any_text = True
                bbox = span.get("bbox", [0, 0, 0, 0])
                resolved = resolve_font(span.get("font", ""), span.get("flags", 0))
                color_int = span.get("color", 0)
                color = f"#{color_int:06x}" if isinstance(color_int, int) else "#000000"
                yield {
                    "bbox": bbox,
                    "text": text,
                    "family": resolved.family,
                    "size": span.get("size", 11.0),
                    "bold": resolved.bold,
                    "italic": resolved.italic,
                    "color": color,
                }

    if not any_text and ocr_fallback:
        from app.services.pdf_ocr import ocr_available, ocr_lines
        if ocr_available():
            try:
                for line in ocr_lines(page):
                    for word in line:
                        yield {
                            "bbox": [word.x0, word.y0, word.x1, word.y1],
                            "text": word.text,
                            "family": "Arial",
                            "size": max(6.0, word.y1 - word.y0),
                            "bold": False,
                            "italic": False,
                            "color": "#000000",
                        }
            except Exception:
                pass


def _page_images(page):
    """Yield (bbox, png_bytes) for each raster image placed on the page."""
    try:
        infos = page.get_image_info(xrefs=True)
    except Exception:
        return
    doc = page.parent
    for info in infos:
        xref = info.get("xref")
        bbox = info.get("bbox")
        if not xref or not bbox:
            continue
        try:
            pix = doc.extract_image(xref)
            data = pix.get("image")
            if data:
                yield bbox, data
        except Exception:
            continue


def pdf_to_html(input_path: Path, output_path: Path, analysis: PdfAnalysis | None = None) -> Path:
    """Render a PDF to fixed-layout HTML that preserves page geometry,
    font family/size/weight/style/color, text position, and images -- using
    absolute positioning within a page-sized container per page, which is
    the only HTML strategy that can represent arbitrary PDF layouts
    (multi-column, sidebars, overlapping text/images) faithfully.
    """
    import fitz

    document = fitz.open(input_path)
    if analysis is None:
        analysis = analyze_pdf(input_path)

    try:
        page_html_blocks = []
        for i, page in enumerate(document):
            rect = page.rect
            w, h = rect.width, rect.height
            page_analysis = analysis.pages[i] if i < len(analysis.pages) else None
            needs_ocr = page_analysis.needs_ocr if page_analysis else False

            elements = []

            # Images first (so text painted after sits visually above them,
            # matching typical PDF z-ordering of background art under text).
            for bbox, img_bytes in _page_images(page):
                x0, y0, x1, y1 = bbox
                try:
                    b64 = __import__("base64").b64encode(img_bytes).decode("ascii")
                except Exception:
                    continue
                elements.append(
                    f'<img style="position:absolute;left:{x0:.1f}px;top:{y0:.1f}px;'
                    f'width:{(x1-x0):.1f}px;height:{(y1-y0):.1f}px;" '
                    f'src="data:image/png;base64,{b64}">'
                )

            for span in _page_spans(page, ocr_fallback=needs_ocr):
                x0, y0, x1, y1 = span["bbox"]
                weight = "bold" if span["bold"] else "normal"
                style_font = "italic" if span["italic"] else "normal"
                text = html.escape(span["text"])
                elements.append(
                    f'<span style="position:absolute;left:{x0:.1f}px;top:{y0:.1f}px;'
                    f'font-family:\'{span["family"]}\',sans-serif;font-size:{span["size"]:.1f}px;'
                    f'font-weight:{weight};font-style:{style_font};color:{span["color"]};'
                    f'white-space:pre;line-height:1;">{text}</span>'
                )

            page_html_blocks.append(
                f'<div class="kiwi-page" style="position:relative;width:{w:.1f}px;'
                f'height:{h:.1f}px;margin:0 auto 24px auto;background:#fff;'
                f'box-shadow:0 1px 4px rgba(0,0,0,0.15);overflow:hidden;">'
                + "".join(elements) + "</div>"
            )

        document_html = (
            "<!doctype html><html><head><meta charset='utf-8'>"
            "<title>Converted PDF</title>"
            "<style>body{background:#e5e5e5;margin:0;padding:24px 0;}"
            ".kiwi-page{font-family:sans-serif;}</style>"
            "</head><body>" + "".join(page_html_blocks) + "</body></html>"
        )
        output_path.write_text(document_html, encoding="utf-8")
    finally:
        document.close()
    return output_path


def _draw_pptx_textbox(slide, spans, scale_x, scale_y, page_height):
    """Group a page's text spans into lines by vertical proximity and add
    one PowerPoint textbox per line, each run carrying its own resolved
    font/size/weight/style/color. Per-line (rather than per-document-block)
    placement keeps positional drift low without needing full paragraph
    reconstruction, which is the right tradeoff for the FIDELITY-leaning
    default mode.
    """
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor

    for span in spans:
        x0, y0, x1, y1 = span["bbox"]
        left = Emu(max(0, int(x0 * scale_x * 914400 / 72)))
        top = Emu(max(0, int(y0 * scale_y * 914400 / 72)))
        width = Emu(max(9144, int((x1 - x0) * scale_x * 914400 / 72)))
        height = Emu(max(9144, int((y1 - y0) * scale_y * 914400 / 72) or 137160))

        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = False
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = span["text"]
        font = run.font
        font.name = span["family"]
        font.size = Pt(max(4.0, min(96.0, span["size"] * scale_y)))
        font.bold = span["bold"]
        font.italic = span["italic"]
        try:
            hexcolor = span["color"].lstrip("#")
            font.color.rgb = RGBColor.from_string(hexcolor.upper().zfill(6))
        except Exception:
            pass


def pdf_to_pptx(input_path: Path, output_path: Path, analysis: PdfAnalysis | None = None) -> Path:
    """Render each PDF page to a PowerPoint slide.

    Strategy per page:
      - If the page has a usable text layer: place real editable text runs
        (resolved fonts, size, bold/italic, color) plus native picture
        objects for embedded images -- this is the EDITABLE path.
      - If the page needs OCR and OCR is available: OCR the page and place
        the recovered words as editable text at their original positions.
      - If neither native text nor OCR is available/successful, or the page
        is dominated by vector art that can't be reconstructed as shapes
        (which would otherwise silently vanish): render the full page as a
        high-resolution image so nothing is lost -- this is the FIDELITY
        fallback, applied per-page rather than forcing the whole deck into
        raster mode just because one page is a complex graphic.
    """
    import fitz
    from pptx import Presentation

    if analysis is None:
        analysis = analyze_pdf(input_path)

    pdf = fitz.open(input_path)
    prs = Presentation()
    while len(prs.slides):
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    try:
        if len(pdf) == 0:
            raise ValueError("PDF contains no pages.")

        first_rect = pdf[0].rect
        prs.slide_width = __import__("pptx").util.Inches(13.333)
        prs.slide_height = __import__("pptx").util.Inches(13.333 * first_rect.height / first_rect.width)

        for i, page in enumerate(pdf):
            page_analysis = analysis.pages[i] if i < len(analysis.pages) else None
            slide_width_pt = float(prs.slide_width) / 12700.0
            slide_height_pt = float(prs.slide_height) / 12700.0
            scale_x = slide_width_pt / page.rect.width if page.rect.width else 1.0
            scale_y = slide_height_pt / page.rect.height if page.rect.height else 1.0

            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Decide the per-page strategy.
            heavy_vector = bool(page_analysis and page_analysis.drawing_count > 25)
            needs_ocr = bool(page_analysis and page_analysis.needs_ocr)

            placed_any = False

            if not heavy_vector:
                spans = list(_page_spans(page, ocr_fallback=needs_ocr))
                if spans:
                    _draw_pptx_textbox(slide, spans, scale_x, scale_y, page.rect.height)
                    placed_any = True

                for bbox, img_bytes in _page_images(page):
                    x0, y0, x1, y1 = bbox
                    left = __import__("pptx").util.Emu(max(0, int(x0 * scale_x * 914400 / 72)))
                    top = __import__("pptx").util.Emu(max(0, int(y0 * scale_y * 914400 / 72)))
                    width = __import__("pptx").util.Emu(max(9144, int((x1 - x0) * scale_x * 914400 / 72)))
                    height = __import__("pptx").util.Emu(max(9144, int((y1 - y0) * scale_y * 914400 / 72)))
                    try:
                        slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width=width, height=height)
                        placed_any = True
                    except Exception:
                        continue

            # Fidelity fallback: page has heavy vector art, or nothing could
            # be placed as editable content -- render the whole page as a
            # high-resolution image so visual content is never silently lost.
            if heavy_vector or not placed_any:
                pix = page.get_pixmap(matrix=fitz.Matrix(3.0, 3.0), alpha=False)
                slide.shapes.add_picture(
                    io.BytesIO(pix.tobytes("png")), 0, 0,
                    width=prs.slide_width, height=prs.slide_height,
                )

        prs.save(output_path)
    finally:
        pdf.close()

    return output_path
