from __future__ import annotations

import html
import io
import re
from pathlib import Path

from app.core.config import settings
from app.core.files import command_path, run_command
from app.core.formats import TEXT_DOCUMENTS, extension_of, valid_document_targets


LIBREOFFICE_FILTER_MAP = {
    "pdf": "pdf",
    "docx": "docx",
    "doc": "doc:MS Word 97",
    "odt": "odt",
    "rtf": "rtf",
    "txt": "txt",
    "html": "html",
    "htm": "html",
    "xlsx": "xlsx",
    "xls": "xls:MS Excel 97",
    "ods": "ods",
    "csv": "csv",
    "pptx": "pptx",
    "ppt": "ppt:MS PowerPoint 97",
    "odp": "odp",
}


def _soffice() -> str:
    return command_path(
        settings.soffice_bin,
        [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ],
    ) or ""


def _strip_markdown(value: str) -> str:
    value = re.sub(r"```.*?```", "", value, flags=re.S)
    value = re.sub(r"`([^`]*)`", r"\1", value)
    value = re.sub(r"!\[[^]]*\]\([^)]*\)", "", value)
    value = re.sub(r"\[([^]]+)\]\([^)]*\)", r"\1", value)
    value = re.sub(r"^#{1,6}\s*", "", value, flags=re.M)
    value = re.sub(r"[*_~]", "", value)
    return value


def _html_to_text(value: str) -> str:
    value = re.sub(r"<br\s*/?>", "\n", value, flags=re.I)
    value = re.sub(r"</(p|div|h[1-6]|li|tr)>", "\n", value, flags=re.I)
    value = re.sub(r"<[^>]+>", "", value)
    return html.unescape(value)


def _pdf_to_text(input_path: Path, output_path: Path) -> Path:
    import fitz

    document = fitz.open(input_path)
    try:
        text = "\n\n".join(page.get_text("text") for page in document)
        output_path.write_text(text, encoding="utf-8")
    finally:
        document.close()
    return output_path


def _pdf_to_html(input_path: Path, output_path: Path) -> Path:
    import fitz

    document = fitz.open(input_path)
    try:
        sections = []
        for index, page in enumerate(document, start=1):
            page_text = html.escape(page.get_text("text"))
            page_text = page_text.replace("\n", "<br>\n")
            sections.append(
                f"<section><h2>Page {index}</h2><p>{page_text}</p></section>"
            )

        document_html = (
            "<!doctype html><html><head><meta charset='utf-8'>"
            "<title>Converted PDF</title></head><body>"
            + "".join(sections)
            + "</body></html>"
        )
        output_path.write_text(document_html, encoding="utf-8")
    finally:
        document.close()
    return output_path


def _pdf_block_text(block: dict) -> str:
    lines = []
    for line in block.get("lines", []):
        spans = [span.get("text", "") for span in line.get("spans", [])]
        text = "".join(spans).strip()
        if text:
            lines.append(text)
    return "\n".join(lines).strip()


def _block_font_size(block: dict) -> float:
    sizes: list[float] = []
    for line in block.get("lines", []):
        for span in line.get("spans", []):
            try:
                size = float(span.get("size", 11))
            except (TypeError, ValueError):
                size = 11
            if size > 0:
                sizes.append(size)
    return max(8.0, min(30.0, max(sizes, default=11.0)))


def _pdf_to_docx(input_path: Path, output_path: Path) -> Path:
    """Create an editable DOCX approximation of a PDF.

    General PDFs do not contain enough semantic information to guarantee a
    perfect Word reconstruction. This converter extracts positioned text
    blocks and embedded raster images while preserving page boundaries.
    """
    import fitz
    from docx import Document
    from docx.enum.text import WD_BREAK, WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt

    pdf = fitz.open(input_path)
    doc = Document()

    section = doc.sections[0]
    section.top_margin = Inches(0.55)
    section.bottom_margin = Inches(0.55)
    section.left_margin = Inches(0.65)
    section.right_margin = Inches(0.65)

    try:
        for page_index, page in enumerate(pdf):
            page_dict = page.get_text("dict")
            text_blocks = [b for b in page_dict.get("blocks", []) if b.get("type") == 0]
            image_blocks = [b for b in page_dict.get("blocks", []) if b.get("type") == 1 and b.get("image")]

            # Reading order from the PDF block coordinates.
            text_blocks.sort(key=lambda block: (block.get("bbox", [0, 0, 0, 0])[1], block.get("bbox", [0, 0, 0, 0])[0]))

            if not text_blocks and not image_blocks:
                paragraph = doc.add_paragraph()
                paragraph.add_run(f"[Page {page_index + 1} contains no extractable text or raster images]")
            else:
                for block in text_blocks:
                    text = _pdf_block_text(block)
                    if not text:
                        continue

                    paragraph = doc.add_paragraph()
                    paragraph.paragraph_format.space_after = Pt(3)
                    paragraph.paragraph_format.line_spacing = 1.0

                    size = _block_font_size(block)
                    first_span = None
                    for line in block.get("lines", []):
                        for span in line.get("spans", []):
                            first_span = span
                            break
                        if first_span:
                            break

                    run = paragraph.add_run(text)
                    run.font.size = Pt(size)
                    run.font.name = "Arial"

                    if first_span and (int(first_span.get("flags", 0)) & 16):
                        run.bold = True

                # Embedded images are inserted after the text for the page.
                # Their original coordinates cannot always be represented by
                # a flow layout in DOCX, so preserving the image itself is the
                # safer behavior than silently dropping it.
                for image_index, block in enumerate(image_blocks, start=1):
                    image_bytes = block.get("image")
                    if not image_bytes:
                        continue
                    paragraph = doc.add_paragraph()
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    try:
                        paragraph.add_run().add_picture(io.BytesIO(image_bytes), width=Inches(6.2))
                    except Exception:
                        continue

            if page_index < len(pdf) - 1:
                doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)

        doc.save(output_path)
    finally:
        pdf.close()

    return output_path


def _pdf_to_pptx(input_path: Path, output_path: Path) -> Path:
    """Create an editable PPTX approximation of a PDF.

    Each PDF page becomes a slide. Text becomes native PowerPoint text boxes
    and embedded raster images become native picture objects. Complex vector
    artwork may not map 1:1, which is inherent in PDF -> PowerPoint conversion.
    """
    import fitz
    from pptx import Presentation
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.util import Inches, Pt

    pdf = fitz.open(input_path)
    prs = Presentation()

    # Remove the default starter slide so every PDF page maps exactly to one slide.
    while len(prs.slides):
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    try:
        if len(pdf) == 0:
            raise ValueError("PDF contains no pages.")

        first_page = pdf[0]
        first_rect = first_page.rect
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(10 * first_rect.height / first_rect.width)

        for page in pdf:
            slide_width_pt = float(prs.slide_width) / 12700.0
            slide_height_pt = float(prs.slide_height) / 12700.0
            scale_x = slide_width_pt / page.rect.width if page.rect.width else 1.0
            scale_y = slide_height_pt / page.rect.height if page.rect.height else 1.0

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            page_dict = page.get_text("dict")
            blocks = page_dict.get("blocks", [])
            blocks.sort(key=lambda block: (block.get("bbox", [0, 0, 0, 0])[1], block.get("bbox", [0, 0, 0, 0])[0]))

            for block in blocks:
                bbox = block.get("bbox") or [0, 0, 0, 0]
                x0, y0, x1, y1 = map(float, bbox)
                left = Inches(max(0.0, x0 * scale_x / 72.0))
                top = Inches(max(0.0, y0 * scale_y / 72.0))
                width = Inches(max(0.05, (x1 - x0) * scale_x / 72.0))
                height = Inches(max(0.05, (y1 - y0) * scale_y / 72.0))

                if block.get("type") == 1 and block.get("image"):
                    try:
                        slide.shapes.add_picture(
                            io.BytesIO(block["image"]),
                            left,
                            top,
                            width=width,
                            height=height,
                        )
                    except Exception:
                        continue
                    continue

                if block.get("type") != 0:
                    continue

                text = _pdf_block_text(block)
                if not text:
                    continue

                shape = slide.shapes.add_textbox(
                    left,
                    top,
                    width,
                    height,
                )

                frame = shape.text_frame
                frame.clear()
                frame.word_wrap = True
                frame.auto_size = MSO_AUTO_SIZE.NONE

                paragraph = frame.paragraphs[0]
                paragraph.text = text

                font = paragraph.runs[0].font
                font.name = "Arial"
                font.size = Pt(max(6, min(32, _block_font_size(block))))

            # If a page had no extractable objects, preserve its visual content
            # as a full-page raster image rather than returning an empty slide.
            if not slide.shapes:
                pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
                slide.shapes.add_picture(
                    io.BytesIO(pix.tobytes("png")),
                    0,
                    0,
                    width=prs.slide_width,
                    height=prs.slide_height,
                )

        prs.save(output_path)
    finally:
        pdf.close()

    return output_path


def convert_with_libreoffice(input_path: Path, target_ext: str, out_dir: Path) -> Path:
    binary = _soffice()
    if not binary:
        raise RuntimeError(
            "LibreOffice is required for this conversion but was not found. "
            "Install LibreOffice, restart Kiwi, or set KIWI_SOFFICE_BIN."
        )

    filter_name = LIBREOFFICE_FILTER_MAP.get(target_ext, target_ext)
    args = [
        binary,
        "--headless",
        "--norestore",
        "--nolockcheck",
        "--nodefault",
        "--convert-to",
        filter_name,
        "--outdir",
        str(out_dir),
        str(input_path),
    ]

    try:
        run_command(args, timeout=240)
    except RuntimeError as exc:
        raise RuntimeError(
            f"LibreOffice could not convert '{input_path.name}' to .{target_ext}. {exc}"
        ) from exc

    produced = out_dir / f"{input_path.stem}.{target_ext}"
    if produced.exists():
        return produced

    matches = [
        path
        for path in out_dir.glob(f"{input_path.stem}.*")
        if path.is_file() and path.suffix.lower().lstrip(".") == target_ext
    ]
    if matches:
        return matches[0]

    raise RuntimeError(
        f"LibreOffice completed without producing the requested .{target_ext} output."
    )


def convert_document(input_path: Path, target_ext: str, out_dir: Path) -> Path:
    source_ext = extension_of(input_path.name)
    target_ext = target_ext.lower().lstrip(".")

    allowed = valid_document_targets(source_ext)
    if target_ext not in allowed:
        raise ValueError(f"Conversion from .{source_ext} to .{target_ext} is not supported.")

    out_dir.mkdir(parents=True, exist_ok=True)
    output_path = out_dir / f"{input_path.stem}.{target_ext}"

    if source_ext == "pdf":
        if target_ext == "txt":
            return _pdf_to_text(input_path, output_path)
        if target_ext == "html":
            return _pdf_to_html(input_path, output_path)
        if target_ext == "docx":
            return _pdf_to_docx(input_path, output_path)
        if target_ext == "pptx":
            return _pdf_to_pptx(input_path, output_path)

    if source_ext in {"md", "markdown"} and target_ext == "txt":
        output_path.write_text(
            _strip_markdown(input_path.read_text(encoding="utf-8", errors="replace")),
            encoding="utf-8",
        )
        return output_path

    if source_ext in {"html", "htm"} and target_ext == "txt":
        output_path.write_text(
            _html_to_text(input_path.read_text(encoding="utf-8", errors="replace")),
            encoding="utf-8",
        )
        return output_path

    if source_ext == "txt" and target_ext == "md":
        output_path.write_text(
            input_path.read_text(encoding="utf-8", errors="replace"),
            encoding="utf-8",
        )
        return output_path

    if source_ext in {"html", "htm"} and target_ext == "md":
        output_path.write_text(
            _html_to_text(input_path.read_text(encoding="utf-8", errors="replace")),
            encoding="utf-8",
        )
        return output_path

    return convert_with_libreoffice(input_path, target_ext, out_dir)
